#!/usr/bin/env python3
import os
import json
import re
import io
import base64
from enum import Enum
from typing import List, Dict, Any, Optional, Union, Tuple
from dataclasses import dataclass
import tempfile
from datetime import datetime

# PowerPoint manipulation
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# NLP
import nltk
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.corpus import stopwords

# Image processing
from PIL import Image

# For MCP integration
from mcp.server.fastmcp import FastMCP, Context
from mcp.server.fastmcp import Image as MCPImage
from mcp.types import Tool, TextContent, ImageContent

# Allow for external mcp instance to be provided
# This will be set when imported by mcp_unified_server.py
_external_mcp = None

# Create local mcp instance (only used when running standalone)
_local_mcp = FastMCP(
    "PowerPoint Tools",
    dependencies=["python-pptx", "nltk", "pillow"]
)

# Helper function to get the correct mcp instance


def get_mcp():
    return _external_mcp if _external_mcp else _local_mcp


# For backward compatibility, expose as mcp
mcp = get_mcp()

# Define PowerPoint Tool operations


class PowerPointTools(str, Enum):
    CREATE_PRESENTATION = "ppt_create_presentation"
    OPEN_PRESENTATION = "ppt_open_presentation"
    SAVE_PRESENTATION = "ppt_save_presentation"
    ADD_SLIDE = "ppt_add_slide"
    ADD_TEXT = "ppt_add_text"
    ADD_IMAGE = "ppt_add_image"
    ADD_CHART = "ppt_add_chart"
    ADD_TABLE = "ppt_add_table"
    ANALYZE_PRESENTATION = "ppt_analyze_presentation"
    GENERATE_PRESENTATION = "ppt_generate_presentation"
    ENHANCE_PRESENTATION = "ppt_enhance_presentation"

# PowerPoint Session Manager


class PowerPointManager:
    """Manages PowerPoint presentation sessions"""

    def __init__(self):
        self.active_presentations = {}
        self.temp_dir = tempfile.mkdtemp()
        self.screenshots = {}  # Add this line to store screenshots

    def store_screenshot(self, name, image_data):
        """Store a screenshot"""
        self.screenshots[name] = image_data
        return f"Screenshot '{name}' stored"

    def create_presentation(self, session_id: str, template_path: Optional[str] = None) -> str:
        """Create a new PowerPoint presentation"""
        try:
            if template_path and os.path.exists(template_path):
                prs = Presentation(template_path)
            else:
                prs = Presentation()

            self.active_presentations[session_id] = {
                "presentation": prs,
                "file_path": None,
                "created_at": datetime.now(),
                "modified_at": datetime.now()
            }

            return f"Created new presentation with session ID: {session_id}"
        except Exception as e:
            return f"Error creating presentation: {str(e)}"

    def open_presentation(self, session_id: str, file_path: str) -> str:
        """Open an existing PowerPoint presentation"""
        try:
            if not os.path.exists(file_path):
                return f"File not found: {file_path}"

            prs = Presentation(file_path)
            self.active_presentations[session_id] = {
                "presentation": prs,
                "file_path": file_path,
                "created_at": datetime.now(),
                "modified_at": datetime.now()
            }

            return f"Opened presentation from {file_path} with session ID: {session_id}"
        except Exception as e:
            return f"Error opening presentation: {str(e)}"

    def save_presentation(self, session_id: str, file_path: Optional[str] = None) -> str:
        """Save the active PowerPoint presentation"""
        try:
            if session_id not in self.active_presentations:
                return f"Session not found: {session_id}"

            prs = self.active_presentations[session_id]["presentation"]
            save_path = file_path or self.active_presentations[session_id]["file_path"]

            if not save_path:
                save_path = os.path.join(self.temp_dir, f"{session_id}.pptx")

            prs.save(save_path)
            self.active_presentations[session_id]["file_path"] = save_path
            self.active_presentations[session_id]["modified_at"] = datetime.now(
            )

            return f"Saved presentation to {save_path}"
        except Exception as e:
            return f"Error saving presentation: {str(e)}"

    def add_slide(self, session_id: str, layout_index: int = 1, title: Optional[str] = None,
                  content: Optional[str] = None) -> str:
        """Add a new slide to the presentation"""
        try:
            if session_id not in self.active_presentations:
                return f"Session not found: {session_id}"

            prs = self.active_presentations[session_id]["presentation"]

            # Get layout from master slides (default to title and content)
            try:
                layout = prs.slide_layouts[layout_index]
            except IndexError:
                layout = prs.slide_layouts[1]  # Default to title and content

            # Add slide
            slide = prs.slides.add_slide(layout)

            # Add title if provided
            if title and hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
                slide.shapes.title.text = title

            # Add content if provided
            if content and len(slide.placeholders) > 1:
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 1:  # Title
                        continue
                    # Content, Text, Body
                    if shape.placeholder_format.type in (2, 3, 7):
                        shape.text = content
                        break

            # Handle empty content placeholders
            for shape in slide.shapes:
                # Check if it's a content placeholder and it's empty
                if (hasattr(shape, 'is_placeholder') and shape.is_placeholder and
                    hasattr(shape, 'placeholder_format') and
                    # Content, Text, Body
                    shape.placeholder_format.type in (2, 3, 7) and
                        hasattr(shape, 'text') and not shape.text.strip()):
                    try:
                        # First try: Remove the element directly
                        sp = shape._element
                        sp.getparent().remove(sp)
                    except:
                        try:
                            # Second try: Hide the shape by setting size to minimal
                            shape.width = 0
                            shape.height = 0
                            # Move it off-slide
                            shape.left = -10000
                            shape.top = -10000
                        except:
                            # If all else fails, just leave it
                            pass

            self.active_presentations[session_id]["modified_at"] = datetime.now(
            )
            return f"Added slide with layout index {layout_index}"
        except Exception as e:
            return f"Error adding slide: {str(e)}"

    def add_text(self, session_id: str, slide_index: int, text: str,
                 left: float = 1.0, top: float = 1.0, width: float = 8.0, height: float = 1.0,
                 font_size: int = 18, font_name: str = 'Calibri', bold: bool = False,
                 italic: bool = False, color: str = '000000') -> str:
        """Add text box to a slide"""
        try:
            if session_id not in self.active_presentations:
                return f"Session not found: {session_id}"

            prs = self.active_presentations[session_id]["presentation"]

            try:
                slide = prs.slides[slide_index]
            except IndexError:
                return f"Slide index {slide_index} out of range"

            # Convert measurements to PowerPoint's units (inches)
            left_inches = Inches(left)
            top_inches = Inches(top)
            width_inches = Inches(width)
            height_inches = Inches(height)

            # Add text box
            textbox = slide.shapes.add_textbox(
                left_inches, top_inches, width_inches, height_inches)
            textframe = textbox.text_frame
            textframe.text = text

            # Format text
            paragraph = textframe.paragraphs[0]
            run = paragraph.runs[0]
            run.font.size = Pt(font_size)
            run.font.name = font_name
            run.font.bold = bold
            run.font.italic = italic

            # Set color (expecting hex color code without #)
            if color and len(color) == 6:
                try:
                    r = int(color[0:2], 16)
                    g = int(color[2:4], 16)
                    b = int(color[4:6], 16)
                    run.font.color.rgb = RGBColor(r, g, b)
                except ValueError:
                    pass  # Invalid color code, ignore

            self.active_presentations[session_id]["modified_at"] = datetime.now(
            )
            return f"Added text box to slide {slide_index}"
        except Exception as e:
            return f"Error adding text: {str(e)}"

    def add_image(self, session_id: str, slide_index: int, image_path: str,
                  left: float = 1.0, top: float = 1.0, width: Optional[float] = None,
                  height: Optional[float] = None) -> str:
        """Add image to a slide"""
        try:
            if session_id not in self.active_presentations:
                return f"Session not found: {session_id}"

            prs = self.active_presentations[session_id]["presentation"]

            try:
                slide = prs.slides[slide_index]
            except IndexError:
                return f"Slide index {slide_index} out of range"

            # Check if image file exists
            if not os.path.exists(image_path):
                return f"Image file not found: {image_path}"

            # Convert measurements to PowerPoint's units (inches)
            left_inches = Inches(left)
            top_inches = Inches(top)

            # Determine image size
            if width is not None and height is not None:
                width_inches = Inches(width)
                height_inches = Inches(height)
                slide.shapes.add_picture(
                    image_path, left_inches, top_inches, width_inches, height_inches)
            else:
                slide.shapes.add_picture(image_path, left_inches, top_inches)

            self.active_presentations[session_id]["modified_at"] = datetime.now(
            )
            return f"Added image to slide {slide_index}"
        except Exception as e:
            return f"Error adding image: {str(e)}"

    def add_chart(self, session_id: str, slide_index: int, chart_type: str,
                  categories: List[str], series_names: List[str], series_values: List[List[float]],
                  left: float = 1.0, top: float = 1.0, width: float = 8.0, height: float = 5.0,
                  chart_title: Optional[str] = None) -> str:
        """Add chart to a slide"""
        try:
            if session_id not in self.active_presentations:
                return f"Session not found: {session_id}"

            prs = self.active_presentations[session_id]["presentation"]

            try:
                slide = prs.slides[slide_index]
            except IndexError:
                return f"Slide index {slide_index} out of range"

            # Convert measurements to PowerPoint's units (inches)
            left_inches = Inches(left)
            top_inches = Inches(top)
            width_inches = Inches(width)
            height_inches = Inches(height)

            # Map chart type string to PowerPoint chart type
            chart_type_map = {
                'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE,
                'area': XL_CHART_TYPE.AREA,
                'scatter': XL_CHART_TYPE.XY_SCATTER,
                'radar': XL_CHART_TYPE.RADAR,
                'stock': XL_CHART_TYPE.STOCK_HLOC,
                'surface': XL_CHART_TYPE.SURFACE,
                'doughnut': XL_CHART_TYPE.DOUGHNUT,
                'bubble': XL_CHART_TYPE.BUBBLE
            }

            xl_chart_type = chart_type_map.get(
                chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)

            # Create chart data
            chart_data = CategoryChartData()
            chart_data.categories = categories

            # Add series
            for i, series_name in enumerate(series_names):
                if i < len(series_values):
                    chart_data.add_series(series_name, series_values[i])

            # Add chart to slide
            chart = slide.shapes.add_chart(xl_chart_type, left_inches, top_inches,
                                           width_inches, height_inches, chart_data).chart

            # Set chart title if provided
            if chart_title:
                chart.has_title = True
                chart.chart_title.text_frame.text = chart_title

            self.active_presentations[session_id]["modified_at"] = datetime.now(
            )
            return f"Added {chart_type} chart to slide {slide_index}"
        except Exception as e:
            return f"Error adding chart: {str(e)}"

    def add_table(self, session_id: str, slide_index: int, rows: int, cols: int,
                  data: List[List[str]], left: float = 1.0, top: float = 1.0,
                  width: float = 8.0, height: float = 5.0) -> str:
        """Add table to a slide"""
        try:
            if session_id not in self.active_presentations:
                return f"Session not found: {session_id}"

            prs = self.active_presentations[session_id]["presentation"]

            try:
                slide = prs.slides[slide_index]
            except IndexError:
                return f"Slide index {slide_index} out of range"

            # Convert measurements to PowerPoint's units (inches)
            left_inches = Inches(left)
            top_inches = Inches(top)
            width_inches = Inches(width)
            height_inches = Inches(height)

            # Add table
            table = slide.shapes.add_table(rows, cols, left_inches, top_inches,
                                           width_inches, height_inches).table

            # Populate table data
            for r in range(min(rows, len(data))):
                for c in range(min(cols, len(data[r]) if r < len(data) else 0)):
                    cell = table.cell(r, c)
                    cell.text = str(data[r][c])

            self.active_presentations[session_id]["modified_at"] = datetime.now(
            )
            return f"Added table with {rows} rows and {cols} columns to slide {slide_index}"
        except Exception as e:
            return f"Error adding table: {str(e)}"

    def analyze_presentation(self, session_id: str) -> Dict[str, Any]:
        """Analyze the content and structure of a presentation"""
        try:
            if session_id not in self.active_presentations:
                return {"error": f"Session not found: {session_id}"}

            prs = self.active_presentations[session_id]["presentation"]

            analysis = {
                "total_slides": len(prs.slides),
                "slides": [],
                "word_count": 0,
                "total_images": 0,
                "total_charts": 0,
                "total_tables": 0,
                "average_words_per_slide": 0,
                "slide_titles": [],
                "presentation_structure": []
            }

            for i, slide in enumerate(prs.slides):
                slide_analysis = self._analyze_slide(slide)
                slide_analysis["slide_number"] = i
                analysis["slides"].append(slide_analysis)

                # Update presentation-wide statistics
                analysis["word_count"] += slide_analysis["word_count"]
                analysis["total_images"] += slide_analysis["image_count"]
                analysis["total_charts"] += slide_analysis["chart_count"]
                analysis["total_tables"] += slide_analysis["table_count"]

                # Get slide title
                if slide_analysis["title"]:
                    analysis["slide_titles"].append(slide_analysis["title"])
                    analysis["presentation_structure"].append({
                        "slide_number": i,
                        "title": slide_analysis["title"]
                    })
                else:
                    analysis["presentation_structure"].append({
                        "slide_number": i,
                        "title": f"Slide {i+1} (No Title)"
                    })

            # Calculate averages
            if analysis["total_slides"] > 0:
                analysis["average_words_per_slide"] = analysis["word_count"] / \
                    analysis["total_slides"]

            return analysis
        except Exception as e:
            return {"error": f"Error analyzing presentation: {str(e)}"}

    def _analyze_slide(self, slide) -> Dict[str, Any]:
        """Analyze a single slide"""
        slide_analysis = {
            "title": "",
            "word_count": 0,
            "text_content": "",
            "image_count": 0,
            "chart_count": 0,
            "table_count": 0,
            "elements": []
        }

        # Get slide title
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title is not None:
            slide_analysis["title"] = slide.shapes.title.text

        # Analyze shapes
        for shape in slide.shapes:
            # Text analysis
            if shape.has_text_frame:
                text = ""
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"

                if text.strip():
                    slide_analysis["text_content"] += text
                    words = text.split()
                    slide_analysis["word_count"] += len(words)

                    slide_analysis["elements"].append({
                        "type": "text",
                        "content": text.strip(),
                        "word_count": len(words)
                    })

            # Image analysis
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_analysis["image_count"] += 1
                slide_analysis["elements"].append({
                    "type": "image"
                })

            # Chart analysis
            if shape.has_chart:
                slide_analysis["chart_count"] += 1
                slide_analysis["elements"].append({
                    "type": "chart",
                    "chart_type": str(shape.chart.chart_type)
                })

            # Table analysis
            if shape.has_table:
                slide_analysis["table_count"] += 1
                rows = shape.table.rows
                cols = shape.table.columns

                table_data = []
                for r in range(len(rows)):
                    row_data = []
                    for c in range(len(cols)):
                        cell_text = shape.table.cell(r, c).text
                        row_data.append(cell_text)
                    table_data.append(row_data)

                slide_analysis["elements"].append({
                    "type": "table",
                    "rows": len(rows),
                    "columns": len(cols),
                    "data": table_data
                })

        return slide_analysis

    def generate_enhancement_suggestions(self, session_id: str) -> Dict[str, Any]:
        """Generate suggestions to enhance the presentation"""
        try:
            analysis = self.analyze_presentation(session_id)

            if "error" in analysis:
                return {"error": analysis["error"]}

            suggestions = {
                "overall_suggestions": [],
                "slide_suggestions": []
            }

            # Overall suggestions
            self._add_overall_suggestions(analysis, suggestions)

            # Per-slide suggestions
            for slide_analysis in analysis["slides"]:
                slide_suggestions = self._generate_slide_suggestions(
                    slide_analysis)
                suggestions["slide_suggestions"].append({
                    "slide_number": slide_analysis["slide_number"],
                    "title": slide_analysis["title"] or f"Slide {slide_analysis['slide_number']+1}",
                    "suggestions": slide_suggestions
                })

            return suggestions
        except Exception as e:
            return {"error": f"Error generating enhancement suggestions: {str(e)}"}

    def _add_overall_suggestions(self, analysis, suggestions):
        """Add overall presentation suggestions"""
        # Check for consistent structure
        if analysis["total_slides"] > 1:
            suggestions["overall_suggestions"].append({
                "type": "structure",
                "suggestion": "Consider adding an agenda or table of contents slide at the beginning"
            })

        # Check word count
        if analysis["average_words_per_slide"] > 100:
            suggestions["overall_suggestions"].append({
                "type": "content_density",
                "suggestion": "Presentation slides have too much text (average over 100 words per slide). Consider breaking content into more slides or reducing text."
            })

        # Check slide count
        if analysis["total_slides"] > 15:
            suggestions["overall_suggestions"].append({
                "type": "length",
                "suggestion": "Presentation is quite long. Consider condensing or breaking into multiple presentations."
            })

        # Check visual elements
        if analysis["total_images"] + analysis["total_charts"] < analysis["total_slides"] / 2:
            suggestions["overall_suggestions"].append({
                "type": "visuals",
                "suggestion": "Add more visual elements (images, charts) to increase engagement."
            })

    def _generate_slide_suggestions(self, slide_analysis):
        """Generate suggestions for a single slide"""
        suggestions = []

        # Check for title
        if not slide_analysis["title"]:
            suggestions.append({
                "type": "structure",
                "suggestion": "Add a clear title to this slide"
            })

        # Check word count on slide
        if slide_analysis["word_count"] > 100:
            suggestions.append({
                "type": "content_density",
                "suggestion": "Slide has too much text. Consider breaking into multiple slides or reducing text."
            })
        elif slide_analysis["word_count"] < 10 and slide_analysis["image_count"] == 0 and slide_analysis["chart_count"] == 0:
            suggestions.append({
                "type": "content",
                "suggestion": "Slide has very little content. Consider adding more information or visuals."
            })

        # Check balance of elements
        if slide_analysis["word_count"] > 50 and slide_analysis["image_count"] + slide_analysis["chart_count"] == 0:
            suggestions.append({
                "type": "balance",
                "suggestion": "Text-heavy slide. Consider adding relevant images or charts to illustrate points."
            })

        return suggestions

    def generate_presentation_from_content(self, session_id: str, title: str, content: Dict[str, Any]) -> str:
        """Generate a presentation from structured content"""
        try:
            # Create new presentation
            result = self.create_presentation(session_id)

            if "Error" in result:
                return result

            prs = self.active_presentations[session_id]["presentation"]

            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            if "subtitle" in content:
                slide.placeholders[1].text = content["subtitle"]

            # Add content slides
            if "slides" in content and isinstance(content["slides"], list):
                for slide_content in content["slides"]:
                    self._add_content_slide(prs, slide_content)

            return f"Generated presentation with title: {title}"
        except Exception as e:
            return f"Error generating presentation: {str(e)}"

    def _add_content_slide(self, prs, slide_content):
        """Add a content slide based on structured data"""
        # Determine slide layout
        layout_index = 1  # Default to Title and Content

        if "layout" in slide_content:
            layout_map = {
                "title": 0,
                "title_content": 1,
                "section": 2,
                "two_content": 3,
                "comparison": 4,
                "title_only": 5,
                "blank": 6,
                "content_caption": 7,
                "picture_caption": 8
            }
            layout_index = layout_map.get(slide_content["layout"].lower(), 1)

        # Try to use specified layout, fallback to simpler layouts if not available
        try:
            layout = prs.slide_layouts[layout_index]
        except IndexError:
            # Fallback to common layouts
            for index in [1, 5, 6]:
                try:
                    layout = prs.slide_layouts[index]
                    break
                except IndexError:
                    continue
            else:
                # If all fallbacks fail, use the first available layout
                layout = prs.slide_layouts[0]

        # Add slide
        slide = prs.slides.add_slide(layout)

        # Add title if provided
        if "title" in slide_content and hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
            slide.shapes.title.text = slide_content["title"]

        # Add text content
        if "content" in slide_content:
            # Try to find content placeholder
            for shape in slide.placeholders:
                # Content, Text, Body
                if shape.placeholder_format.type in (2, 3, 7):
                    shape.text = slide_content["content"]
                    break
            else:
                # If no suitable placeholder found, add textbox
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(4)

                textbox = slide.shapes.add_textbox(left, top, width, height)
                textbox.text_frame.text = slide_content["content"]

        # Add image if specified
        if "image_path" in slide_content and os.path.exists(slide_content["image_path"]):
            # Default position
            left = Inches(3)
            top = Inches(3)

            # Try to find image placeholder
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 18:  # Picture
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    slide.shapes.add_picture(
                        slide_content["image_path"], left, top, width, height)
                    break
            else:
                # If no picture placeholder, add image to slide
                slide.shapes.add_picture(
                    slide_content["image_path"], left, top)

        # Add chart if specified
        if "chart" in slide_content:
            chart_info = slide_content["chart"]

            if ("type" in chart_info and "categories" in chart_info and
                    "series_names" in chart_info and "series_values" in chart_info):

                # Default position
                left = Inches(1.5)
                top = Inches(2)
                width = Inches(7)
                height = Inches(5)

                # Try to find chart placeholder
                for shape in slide.placeholders:
                    # Content, Object
                    if shape.placeholder_format.type in (2, 3):
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        break

                # Map chart type
                chart_type_map = {
                    'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                    'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                    'line': XL_CHART_TYPE.LINE,
                    'pie': XL_CHART_TYPE.PIE,
                    'area': XL_CHART_TYPE.AREA
                }

                xl_chart_type = chart_type_map.get(
                    chart_info["type"].lower(),
                    XL_CHART_TYPE.COLUMN_CLUSTERED
                )

                # Create chart data
                chart_data = CategoryChartData()
                chart_data.categories = chart_info["categories"]

                # Add series
                for i, series_name in enumerate(chart_info["series_names"]):
                    if i < len(chart_info["series_values"]):
                        chart_data.add_series(
                            series_name, chart_info["series_values"][i])

                # Add chart to slide
                chart = slide.shapes.add_chart(
                    xl_chart_type, left, top, width, height, chart_data
                ).chart

                # Set chart title if provided
                if "title" in chart_info:
                    chart.has_title = True
                    chart.chart_title.text_frame.text = chart_info["title"]

        # Add table if specified
        if "table" in slide_content:
            table_info = slide_content["table"]

            if "rows" in table_info and "cols" in table_info and "data" in table_info:
                # Default position
                left = Inches(1)
                top = Inches(2.5)
                width = Inches(8)
                height = Inches(4)

                # Try to find table placeholder
                for shape in slide.placeholders:
                    # Content, Object, Text
                    if shape.placeholder_format.type in (2, 3, 7):
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        break

                # Add table
                rows, cols = table_info["rows"], table_info["cols"]
                table = slide.shapes.add_table(
                    rows, cols, left, top, width, height
                ).table

                # Add data
                for r in range(min(rows, len(table_info["data"]))):
                    for c in range(min(cols, len(table_info["data"][r]) if r < len(table_info["data"]) else 0)):
                        cell = table.cell(r, c)
                        cell.text = str(table_info["data"][r][c])

# NLP utilities for PowerPoint content generation


class PowerPointNLP:
    """Natural language processing utilities for PowerPoint content"""

    def __init__(self):
        # Download required NLTK resources if not already available
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt')

        try:
            nltk.data.find('corpora/stopwords')
        except LookupError:
            nltk.download('stopwords')

    def extract_slide_structure(self, text: str) -> List[Dict[str, Any]]:
        """Extract slide structure from plain text content"""
        # Split text into paragraphs
        paragraphs = text.split('\n\n')

        # Extract title (first non-empty paragraph)
        title = "New Presentation"
        content_start = 0

        for i, paragraph in enumerate(paragraphs):
            if paragraph.strip():
                title = paragraph.strip()
                content_start = i + 1
                break

        # Initialize slides
        slides = []

        # Generate title slide
        slides.append({
            "layout": "title",
            "title": title
        })

        # Process remaining paragraphs for content slides
        current_slide = None

        for paragraph in paragraphs[content_start:]:
            paragraph = paragraph.strip()
            if not paragraph:
                continue

            # Check if this is a heading (potential slide title)
            sentences = sent_tokenize(paragraph)
            first_sentence = sentences[0] if sentences else ""

            is_heading = (
                len(first_sentence) < 100 and
                len(sentences) <= 2 and
                not first_sentence.endswith('.')
            )

            if is_heading:
                # Save previous slide if exists
                if current_slide:
                    slides.append(current_slide)

                # Start new slide
                current_slide = {
                    "layout": "title_content",
                    "title": paragraph,
                    "content": ""
                }
            else:
                # Add content to current slide or create new one if needed
                if not current_slide:
                    current_slide = {
                        "layout": "title_content",
                        "title": "Content Slide",
                        "content": paragraph
                    }
                else:
                    current_slide["content"] += paragraph + "\n\n"

        # Add the last slide if not empty
        if current_slide:
            slides.append(current_slide)

        return slides

    def extract_structured_content(self, text: str) -> Dict[str, Any]:
        """Convert unstructured text to structured content for presentation"""
        slides = self.extract_slide_structure(text)

        # Create presentation structure
        presentation = {
            "title": slides[0]["title"] if slides else "New Presentation",
            "slides": slides[1:] if slides else []
        }

        return presentation

    def suggest_visuals(self, text: str) -> List[Dict[str, str]]:
        """Suggest visual elements based on text content"""
        suggestions = []
        sentences = sent_tokenize(text)

        # Check for numerical data that could be visualized
        numbers_pattern = r'\b\d+\b'
        numbers_count = len(re.findall(numbers_pattern, text))

        if numbers_count > 5:
            suggestions.append({
                "type": "chart",
                "suggestion": "Consider adding a chart to visualize the numerical data"
            })

        # Check for lists that could be tables
        list_markers = ['•', '-', '*', '1.', '2.', '3.']
        list_items = 0

        for sentence in sentences:
            for marker in list_markers:
                if sentence.strip().startswith(marker):
                    list_items += 1
                    break

        if list_items >= 3:
            suggestions.append({
                "type": "table",
                "suggestion": "Consider converting the list into a table for better organization"
            })

        # Check for comparison language
        comparison_words = ['versus', 'vs', 'compared to',
                            'better than', 'worse than', 'difference between']
        for word in comparison_words:
            if word in text.lower():
                suggestions.append({
                    "type": "comparison",
                    "suggestion": "Consider using a comparison slide layout with two columns"
                })
                break

        return suggestions


# Initialize PowerPoint Manager
ppt_manager = PowerPointManager()
ppt_nlp = PowerPointNLP()

# MCP Tool implementations


@get_mcp().tool(name=PowerPointTools.CREATE_PRESENTATION)
def ppt_create_presentation(session_id: str, template_path: Optional[str] = None, ctx: Context = None) -> str:
    """Create a new PowerPoint presentation"""
    return ppt_manager.create_presentation(session_id, template_path)


@get_mcp().tool(name=PowerPointTools.OPEN_PRESENTATION)
def ppt_open_presentation(session_id: str, file_path: str, ctx: Context = None) -> str:
    """Open an existing PowerPoint presentation"""
    return ppt_manager.open_presentation(session_id, file_path)


@get_mcp().tool(name=PowerPointTools.SAVE_PRESENTATION)
def ppt_save_presentation(session_id: str, file_path: Optional[str] = None, ctx: Context = None) -> str:
    """Save the active PowerPoint presentation"""
    return ppt_manager.save_presentation(session_id, file_path)


@get_mcp().tool(name=PowerPointTools.ADD_SLIDE)
def ppt_add_slide(session_id: str, layout_index: int = 1, title: Optional[str] = None,
                  content: Optional[str] = None, ctx: Context = None) -> str:
    """Add a new slide to the presentation"""
    return ppt_manager.add_slide(session_id, layout_index, title, content)


@get_mcp().tool(name=PowerPointTools.ADD_TEXT)
def ppt_add_text(session_id: str, slide_index: int, text: str,
                 left: float = 1.0, top: float = 1.0, width: float = 8.0, height: float = 1.0,
                 font_size: int = 18, font_name: str = 'Calibri', bold: bool = False,
                 italic: bool = False, color: str = '000000', ctx: Context = None) -> str:
    """Add text box to a slide"""
    return ppt_manager.add_text(session_id, slide_index, text, left, top, width, height,
                                font_size, font_name, bold, italic, color)


@get_mcp().tool(name=PowerPointTools.ADD_IMAGE)
def ppt_add_image(session_id: str, slide_index: int, image_path: str,
                  left: float = 1.0, top: float = 1.0, width: Optional[float] = None,
                  height: Optional[float] = None, ctx: Context = None) -> str:
    """Add image to a slide"""
    return ppt_manager.add_image(session_id, slide_index, image_path, left, top, width, height)


@get_mcp().tool(name=PowerPointTools.ADD_CHART)
def ppt_add_chart(session_id: str, slide_index: int, chart_type: str,
                  categories: List[str], series_names: List[str], series_values: List[List[float]],
                  left: float = 1.0, top: float = 1.0, width: float = 8.0, height: float = 5.0,
                  chart_title: Optional[str] = None, ctx: Context = None) -> str:
    """Add chart to a slide"""
    return ppt_manager.add_chart(session_id, slide_index, chart_type, categories, series_names,
                                 series_values, left, top, width, height, chart_title)


@get_mcp().tool(name=PowerPointTools.ADD_TABLE)
def ppt_add_table(session_id: str, slide_index: int, rows: int, cols: int,
                  data: List[List[str]], left: float = 1.0, top: float = 1.0,
                  width: float = 8.0, height: float = 5.0, ctx: Context = None) -> str:
    """Add table to a slide"""
    return ppt_manager.add_table(session_id, slide_index, rows, cols, data, left, top, width, height)


@get_mcp().tool(name=PowerPointTools.ANALYZE_PRESENTATION)
def ppt_analyze_presentation(session_id: str, ctx: Context = None) -> str:
    """Analyze the content and structure of a presentation"""
    analysis = ppt_manager.analyze_presentation(session_id)

    if "error" in analysis:
        return analysis["error"]

    return json.dumps(analysis, indent=2)


@get_mcp().tool(name=PowerPointTools.ENHANCE_PRESENTATION)
def ppt_enhance_presentation(session_id: str, ctx: Context = None) -> str:
    """Provide suggestions to enhance the presentation"""
    suggestions = ppt_manager.generate_enhancement_suggestions(session_id)

    if "error" in suggestions:
        return suggestions["error"]

    return json.dumps(suggestions, indent=2)


@get_mcp().tool(name=PowerPointTools.GENERATE_PRESENTATION)
def ppt_generate_presentation(session_id: str, title: str, content: str, ctx: Context = None) -> str:
    """Generate a presentation from text content"""
    try:
        # Process content with NLP
        structured_content = ppt_nlp.extract_structured_content(content)

        # Generate presentation
        return ppt_manager.generate_presentation_from_content(session_id, title, structured_content)
    except Exception as e:
        return f"Error generating presentation: {str(e)}"


@get_mcp().resource("ppt://screenshots/{name}")
def get_ppt_screenshot(name: str) -> Union[MCPImage, str]:
    """Get a PowerPoint screenshot by name"""
    if name in ppt_manager.screenshots:
        return MCPImage(
            data=ppt_manager.screenshots[name],
            format="png"
        )
    return f"Screenshot '{name}' not found"


@get_mcp().resource("ppt://presentations")
def get_ppt_presentations() -> str:
    """Get list of active PowerPoint presentations"""
    presentations = []
    for session_id, session in ppt_manager.active_presentations.items():
        presentations.append({
            "session_id": session_id,
            "file_path": session.get("file_path", "Unsaved"),
            "created_at": session.get("created_at", "").isoformat() if hasattr(session.get("created_at", ""), "isoformat") else "",
            "slides_count": len(session.get("presentation", {}).slides) if hasattr(session.get("presentation", {}), "slides") else 0
        })
    return json.dumps(presentations, indent=2)

# Natural language command processor


class PowerPointCommander:
    """Process natural language commands for PowerPoint operations"""

    def __init__(self, ppt_manager, ppt_nlp):
        self.ppt_manager = ppt_manager
        self.ppt_nlp = ppt_nlp
        self.session_id = None

    def process_command(self, command: str) -> str:
        """Process a natural language command"""
        command = command.lower().strip()

        # Create a new presentation
        if re.search(r'create|new|start', command) and re.search(r'presentation|slide deck|deck|ppt', command):
            if not self.session_id:
                self.session_id = f"session_{datetime.now().strftime('%Y%m%d%H%M%S')}"

            template_path = None
            template_match = re.search(
                r'template\s+(?:from|at|is|=)\s+([\w\./\\]+)', command)
            if template_match:
                template_path = template_match.group(1)

            return self.ppt_manager.create_presentation(self.session_id, template_path)

        # Open a presentation
        elif re.search(r'open|load', command) and re.search(r'presentation|slide deck|deck|ppt', command):
            file_match = re.search(
                r'(?:file|path)\s+(?:is|=)\s+([\w\./\\]+)', command)
            if file_match:
                file_path = file_match.group(1)

                if not self.session_id:
                    self.session_id = f"session_{datetime.now().strftime('%Y%m%d%H%M%S')}"

                return self.ppt_manager.open_presentation(self.session_id, file_path)
            else:
                return "Please specify the file path to open"

        # Save a presentation
        elif re.search(r'save', command) and re.search(r'presentation|slide deck|deck|ppt', command):
            if not self.session_id:
                return "No active presentation session"

            file_path = None
            file_match = re.search(
                r'(?:file|path|to)\s+(?:is|=|as)\s+([\w\./\\]+)', command)
            if file_match:
                file_path = file_match.group(1)

            return self.ppt_manager.save_presentation(self.session_id, file_path)

        # Add a slide
        elif re.search(r'add', command) and re.search(r'slide', command):
            if not self.session_id:
                return "No active presentation session"

            title = None
            title_match = re.search(r'title\s+(?:is|=)\s+([\w\s]+)', command)
            if title_match:
                title = title_match.group(1)

            content = None
            content_match = re.search(
                r'content\s+(?:is|=)\s+([\w\s]+)', command)
            if content_match:
                content = content_match.group(1)

            layout_index = 1
            layout_match = re.search(r'layout\s+(?:is|=)\s+(\d+)', command)
            if layout_match:
                layout_index = int(layout_match.group(1))

            return self.ppt_manager.add_slide(self.session_id, layout_index, title, content)

        # Generate a presentation
        elif re.search(r'generate|create|make', command) and re.search(r'from|with|using', command) and re.search(r'content|text', command):
            if not self.session_id:
                self.session_id = f"session_{datetime.now().strftime('%Y%m%d%H%M%S')}"

            title = "Generated Presentation"
            title_match = re.search(r'title\s+(?:is|=)\s+([\w\s]+)', command)
            if title_match:
                title = title_match.group(1)

            content = "Sample content"
            content_match = re.search(
                r'content\s+(?:is|=)\s+([\w\s]+)', command)
            if content_match:
                content = content_match.group(1)

            structured_content = self.ppt_nlp.extract_structured_content(
                content)
            return self.ppt_manager.generate_presentation_from_content(self.session_id, title, structured_content)

        # Analyze a presentation
        elif re.search(r'analyze|review', command) and re.search(r'presentation|slide deck|deck|ppt', command):
            if not self.session_id:
                return "No active presentation session"

            analysis = self.ppt_manager.analyze_presentation(self.session_id)
            if "error" in analysis:
                return analysis["error"]

            return json.dumps(analysis, indent=2)

        # Enhance a presentation
        elif re.search(r'enhance|improve|suggest', command) and re.search(r'presentation|slide deck|deck|ppt', command):
            if not self.session_id:
                return "No active presentation session"

            suggestions = self.ppt_manager.generate_enhancement_suggestions(
                self.session_id)
            if "error" in suggestions:
                return suggestions["error"]

            return json.dumps(suggestions, indent=2)

        else:
            return "Unrecognized command. Try something like 'create new presentation', 'add slide', or 'analyze presentation'."


# Initialize the commander
ppt_commander = PowerPointCommander(ppt_manager, ppt_nlp)

# Add a natural language command processor tool


@get_mcp().tool(name="ppt_command")
def ppt_command(command: str, ctx: Context = None) -> str:
    """Process a natural language command for PowerPoint operations"""
    return ppt_commander.process_command(command)


def get_ppt_tools():
    """Return the PowerPoint tools for registration with another MCP instance"""
    return {
        PowerPointTools.CREATE_PRESENTATION: ppt_create_presentation,
        PowerPointTools.OPEN_PRESENTATION: ppt_open_presentation,
        PowerPointTools.SAVE_PRESENTATION: ppt_save_presentation,
        PowerPointTools.ADD_SLIDE: ppt_add_slide,
        PowerPointTools.ADD_TEXT: ppt_add_text,
        PowerPointTools.ADD_IMAGE: ppt_add_image,
        PowerPointTools.ADD_CHART: ppt_add_chart,
        PowerPointTools.ADD_TABLE: ppt_add_table,
        PowerPointTools.ANALYZE_PRESENTATION: ppt_analyze_presentation,
        PowerPointTools.ENHANCE_PRESENTATION: ppt_enhance_presentation,
        PowerPointTools.GENERATE_PRESENTATION: ppt_generate_presentation,
        "ppt_command": ppt_command
    }


def set_external_mcp(external_mcp_instance):
    """Set an external MCP instance to use instead of the local one"""
    global _external_mcp
    global mcp
    _external_mcp = external_mcp_instance
    mcp = get_mcp()

    return _external_mcp


# Main execution as MCP Tool
if __name__ == "__main__":
    # Use configuration from environment variables if available
    host = os.environ.get("MCP_HOST", "0.0.0.0")
    port = int(os.environ.get("MCP_PORT", "8000"))
    log_level = os.environ.get("MCP_LOG_LEVEL", "debug")

    # Update configuration
    mcp.config = {
        "host": host,
        "port": port,
        "log_level": log_level
    }

    # Run the server with configuration from mcp.config
    mcp.run()
